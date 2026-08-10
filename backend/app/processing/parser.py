"""Extract raw text from supported document formats."""
from __future__ import annotations

import io
from pathlib import Path

from docx import Document as DocxDocument
from pypdf import PdfReader
from pypdf.errors import PdfReadError

_SUPPORTED_EXTENSIONS = frozenset({
    ".pdf",
    ".docx",
    ".doc",
    ".txt",
    ".md",
    ".markdown",
    ".csv",
    ".xlsx",
    ".pptx",
    ".json",
    ".log",
    ".html",
    ".htm",
})


class ParsingError(Exception):
    """Base class for document parsing failures."""


class UnsupportedFormatError(ParsingError):
    """The file extension is not supported by the parser."""


class CorruptedFileError(ParsingError):
    """The file is recognized but cannot be read (corrupt or invalid content)."""


def parse_file(content: bytes, filename: str, mime_type: str | None = None) -> str:
    """Extract raw text from an uploaded file."""
    extension = Path(filename).suffix.lower()
    if extension not in _SUPPORTED_EXTENSIONS:
        raise UnsupportedFormatError(
            f"Unsupported file extension {extension!r}. "
            f"Supported: {', '.join(sorted(_SUPPORTED_EXTENSIONS))}."
        )

    if not content:
        raise CorruptedFileError(f"File {filename!r} is empty and cannot be parsed.")

    try:
        if extension == ".pdf":
            return _parse_pdf(content, filename)
        if extension in (".docx", ".doc"):
            return _parse_docx(content, filename)
        if extension == ".xlsx":
            return _parse_xlsx(content, filename)
        if extension == ".pptx":
            return _parse_pptx(content, filename)
        return _parse_plain_text(content, filename)
    except ParsingError:
        raise
    except Exception as exc:
        raise CorruptedFileError(f"Failed to parse {filename!r}: {exc}") from exc


def _parse_pdf(content: bytes, filename: str) -> str:
    parts: list[str] = []

    # 1. Primary: PyMuPDF (fitz) if installed for rich text & layout extraction
    try:
        import importlib

        fitz = importlib.import_module("fitz")
        pdf_doc = fitz.open(stream=content, filetype="pdf")
        for page in pdf_doc:
            page_text = page.get_text()
            if page_text and page_text.strip():
                parts.append(page_text.strip())
        if parts:
            return "\n\n".join(parts)
    except Exception:
        parts = []

    # 2. Secondary fallback: pypdf
    try:
        reader = PdfReader(io.BytesIO(content))
    except PdfReadError as exc:
        raise CorruptedFileError(f"PDF file {filename!r} is corrupted or invalid: {exc}") from exc

    if len(reader.pages) == 0:
        raise CorruptedFileError(f"PDF file {filename!r} contains no pages.")

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text and page_text.strip():
            parts.append(page_text.strip())

    if not parts:
        raise CorruptedFileError(f"PDF file {filename!r} contains no extractable text.")

    return "\n\n".join(parts)


def _parse_docx(content: bytes, filename: str) -> str:
    try:
        document = DocxDocument(io.BytesIO(content))
    except Exception as exc:
        raise CorruptedFileError(f"DOCX file {filename!r} is corrupted or invalid: {exc}") from exc

    parts: list[str] = []

    # 1. Extract paragraphs
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    # 2. Extract table cell contents
    for table in document.tables:
        for row in table.rows:
            row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_cells:
                # Deduplicate identical adjacent cell text from merged cells
                deduped: list[str] = []
                for cell_text in row_cells:
                    if not deduped or deduped[-1] != cell_text:
                        deduped.append(cell_text)
                parts.append(" | ".join(deduped))

    if not parts:
        raise CorruptedFileError(f"DOCX file {filename!r} contains no extractable text.")

    extracted_text = "\n".join(parts).strip()
    if not extracted_text:
        raise CorruptedFileError(f"DOCX file {filename!r} resulted in 0 extracted characters.")

    return extracted_text


def _parse_xlsx(content: bytes, filename: str) -> str:
    try:
        import openpyxl

        workbook = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        parts: list[str] = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            parts.append(f"--- Sheet: {sheet_name} ---")
            for row in sheet.iter_rows(values_only=True):
                row_str = " | ".join(str(cell).strip() for cell in row if cell is not None and str(cell).strip())
                if row_str:
                    parts.append(row_str)
        if not parts:
            raise CorruptedFileError(f"XLSX file {filename!r} contains no extractable data.")
        return "\n".join(parts)
    except Exception as exc:
        raise CorruptedFileError(f"XLSX file {filename!r} is corrupted or invalid: {exc}") from exc


def _parse_pptx(content: bytes, filename: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        parts: list[str] = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slide {slide_idx} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    parts.append(shape.text.strip())
        if not parts:
            raise CorruptedFileError(f"PPTX file {filename!r} contains no extractable text.")
        return "\n".join(parts)
    except Exception as exc:
        raise CorruptedFileError(f"PPTX file {filename!r} is corrupted or invalid: {exc}") from exc


def _parse_plain_text(content: bytes, filename: str) -> str:
    for encoding in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            text = content.decode(encoding)
            if text.strip():
                return text
            raise CorruptedFileError(f"Text file {filename!r} contains no extractable text.")
        except UnicodeDecodeError:
            continue

    raise CorruptedFileError(f"Text file {filename!r} uses an unsupported text encoding.")
