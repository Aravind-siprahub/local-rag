"""Structured document parser with hierarchy preservation.

Parser priority (best available):
  1. Docling (optional)
  2. MarkItDown (optional)
  3. Format-specific fallbacks: PyMuPDF, python-docx, openpyxl, python-pptx
"""
from __future__ import annotations

import io
import logging
import re
import uuid
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Iterator

from app.services.metadata import BlockType, DocumentBlock, ParsedDocument
from app.processing.parser import CorruptedFileError, ParsingError, UnsupportedFormatError

logger = logging.getLogger(__name__)

_SUPPORTED_EXTENSIONS = frozenset({
    ".pdf", ".docx", ".doc", ".txt", ".md", ".markdown",
    ".xlsx", ".pptx", ".html", ".htm",
})

_PAGE_NUMBER_RE = re.compile(r"^\s*\d{1,4}\s*$", re.MULTILINE)
_HEADER_FOOTER_RE = re.compile(
    r"^(?:page\s+\d+\s+of\s+\d+|confidential|draft|copyright\s+©).*$",
    re.IGNORECASE | re.MULTILINE,
)
_WATERMARK_RE = re.compile(r"\b(?:DRAFT|CONFIDENTIAL|DO NOT DISTRIBUTE)\b", re.IGNORECASE)

_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+)$", re.MULTILINE)
_MD_CODE_FENCE_RE = re.compile(r"```[\w]*\n([\s\S]*?)```", re.MULTILINE)
_MD_TABLE_ROW_RE = re.compile(r"^\|.+\|\s*$", re.MULTILINE)
_MD_LIST_RE = re.compile(r"^(\s*[-*+]|\s*\d+\.)\s+.+$", re.MULTILINE)
_FAQ_Q_RE = re.compile(
    r"^(?:Q(?:uestion)?[:\.]?\s*|Q\d+[:\.]?\s*)(.+)$",
    re.IGNORECASE | re.MULTILINE,
)
_FAQ_A_RE = re.compile(
    r"^(?:A(?:nswer)?[:\.]?\s*|A\d+[:\.]?\s*)(.+)$",
    re.IGNORECASE | re.MULTILINE,
)
_PY_DEF_RE = re.compile(r"^(def\s+\w+\s*\([^)]*\)\s*(?:->\s*[^:]+)?:)", re.MULTILINE)
_PY_CLASS_RE = re.compile(r"^(class\s+\w+[^:]*:)", re.MULTILINE)


class DocumentParser:
    """Parse uploaded documents into structured blocks with hierarchy."""

    def __init__(self) -> None:
        self._docling_available = self._check_docling()
        self._markitdown_available = self._check_markitdown()

    @staticmethod
    def _check_docling() -> bool:
        try:
            import importlib
            importlib.import_module("docling.document_converter")
            return True
        except ImportError:
            return False

    @staticmethod
    def _check_markitdown() -> bool:
        try:
            import importlib
            importlib.import_module("markitdown")
            return True
        except ImportError:
            return False

    async def parse(
        self,
        content: bytes,
        filename: str,
        document_id: uuid.UUID,
        mime_type: str | None = None,
    ) -> ParsedDocument:
        """Parse document bytes into a structured ParsedDocument."""
        return await __import__("asyncio").to_thread(
            self.parse_sync, content, filename, document_id, mime_type
        )

    def parse_sync(
        self,
        content: bytes,
        filename: str,
        document_id: uuid.UUID,
        mime_type: str | None = None,
    ) -> ParsedDocument:
        """Synchronous parse — suitable for thread-pool execution."""
        extension = Path(filename).suffix.lower()
        if extension not in _SUPPORTED_EXTENSIONS:
            raise UnsupportedFormatError(
                f"Unsupported extension {extension!r}. "
                f"Supported: {', '.join(sorted(_SUPPORTED_EXTENSIONS))}."
            )
        if not content:
            raise CorruptedFileError(f"File {filename!r} is empty.")

        parser_used = "fallback"
        blocks: list[DocumentBlock] = []
        page_count = 0
        language = "en"

        # 1. Docling (optional, best structure preservation).
        if self._docling_available:
            try:
                blocks, page_count, language = self._parse_with_docling(content, filename)
                parser_used = "docling"
            except Exception as exc:
                logger.debug("Docling parse failed for %s: %s", filename, exc)
                blocks = []

        # 2. MarkItDown (optional).
        if not blocks and self._markitdown_available:
            try:
                blocks, page_count, language = self._parse_with_markitdown(content, filename)
                parser_used = "markitdown"
            except Exception as exc:
                logger.debug("MarkItDown parse failed for %s: %s", filename, exc)
                blocks = []

        # 3. Format-specific fallbacks.
        if not blocks:
            blocks, page_count, language, parser_used = self._parse_fallback(
                content, filename, extension
            )

        if not blocks:
            raise CorruptedFileError(f"No extractable content from {filename!r}.")

        cleaned_blocks = self._remove_noise(blocks)
        if not cleaned_blocks:
            raise CorruptedFileError(f"All content filtered as noise from {filename!r}.")

        return ParsedDocument(
            document_id=document_id,
            document_name=filename,
            language=language,
            page_count=page_count,
            blocks=cleaned_blocks,
            source_format=extension.lstrip("."),
            parser_used=parser_used,
        )

    def _parse_with_docling(
        self, content: bytes, filename: str
    ) -> tuple[list[DocumentBlock], int, str]:
        from docling.document_converter import DocumentConverter

        converter = DocumentConverter()
        result = converter.convert(io.BytesIO(content))
        doc = result.document
        blocks: list[DocumentBlock] = []
        page_count = len(getattr(doc, "pages", []) or [])

        for item in doc.iterate_items():
            label = str(getattr(item, "label", "") or "").lower()
            text = str(getattr(item, "text", "") or "").strip()
            if not text:
                continue
            page = getattr(item, "page_no", None) or getattr(item, "page_number", None)

            if "title" in label or "heading" in label or "section" in label:
                level = int(getattr(item, "level", 1) or 1)
                block_type = BlockType.HEADING if level <= 1 else BlockType.SUBHEADING
                blocks.append(DocumentBlock(
                    block_type=block_type, text=text, level=level, page_number=page
                ))
            elif "table" in label:
                blocks.append(DocumentBlock(
                    block_type=BlockType.TABLE, text=text, page_number=page
                ))
            elif "list" in label:
                blocks.append(DocumentBlock(
                    block_type=BlockType.LIST, text=text, page_number=page
                ))
            elif "code" in label:
                blocks.append(DocumentBlock(
                    block_type=BlockType.CODE, text=text, page_number=page
                ))
            else:
                blocks.append(DocumentBlock(
                    block_type=BlockType.PARAGRAPH, text=text, page_number=page
                ))

        return blocks, page_count, "en"

    def _parse_with_markitdown(
        self, content: bytes, filename: str
    ) -> tuple[list[DocumentBlock], int, str]:
        from markitdown import MarkItDown

        md = MarkItDown()
        result = md.convert_stream(io.BytesIO(content), file_extension=Path(filename).suffix)
        markdown_text = str(result.text_content or "")
        blocks = self._parse_markdown_blocks(markdown_text)
        return blocks, 0, self._detect_language(markdown_text)

    def _parse_fallback(
        self, content: bytes, filename: str, extension: str
    ) -> tuple[list[DocumentBlock], int, str, str]:
        if extension == ".pdf":
            return (*self._parse_pdf(content, filename), "pymupdf")
        if extension in (".docx", ".doc"):
            return (*self._parse_docx(content, filename), "python-docx")
        if extension == ".xlsx":
            return (*self._parse_xlsx(content, filename), "openpyxl")
        if extension == ".pptx":
            return (*self._parse_pptx(content, filename), "python-pptx")
        if extension in (".html", ".htm"):
            return (*self._parse_html(content, filename), "html-parser")
        if extension in (".md", ".markdown"):
            text = self._decode_text(content, filename)
            return self._parse_markdown_blocks(text), 0, self._detect_language(text), "markdown"
        # Plain text / default.
        text = self._decode_text(content, filename)
        return self._parse_plain_text_blocks(text), 0, self._detect_language(text), "plain-text"

    def _parse_pdf(
        self, content: bytes, filename: str
    ) -> tuple[list[DocumentBlock], int, str]:
        blocks: list[DocumentBlock] = []
        page_count = 0
        full_text_parts: list[str] = []

        try:
            import fitz  # PyMuPDF

            pdf_doc = fitz.open(stream=content, filetype="pdf")
            page_count = len(pdf_doc)

            for page_idx, page in enumerate(pdf_doc, start=1):
                # Use dict extraction for layout-aware blocks.
                page_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)
                page_blocks = self._extract_pdf_page_blocks(page_dict, page_idx)
                if page_blocks:
                    blocks.extend(page_blocks)
                else:
                    # Fallback to plain text per page.
                    page_text = page.get_text().strip()
                    if page_text:
                        full_text_parts.append(page_text)
                        for para in self._split_paragraphs(page_text):
                            blocks.append(DocumentBlock(
                                block_type=BlockType.PARAGRAPH,
                                text=para,
                                page_number=page_idx,
                            ))

            if blocks:
                combined = "\n\n".join(b.text for b in blocks)
                return blocks, page_count, self._detect_language(combined)
        except Exception as exc:
            logger.debug("PyMuPDF structured parse failed: %s", exc)

        # pypdf fallback.
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError

        try:
            reader = PdfReader(io.BytesIO(content))
        except PdfReadError as exc:
            raise CorruptedFileError(f"PDF {filename!r} is corrupted: {exc}") from exc

        page_count = len(reader.pages)
        if page_count == 0:
            raise CorruptedFileError(f"PDF {filename!r} has no pages.")

        for page_idx, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if not page_text:
                continue
            for para in self._split_paragraphs(page_text):
                blocks.append(DocumentBlock(
                    block_type=BlockType.PARAGRAPH, text=para, page_number=page_idx
                ))

        if not blocks:
            raise CorruptedFileError(f"PDF {filename!r} has no extractable text.")

        combined = "\n\n".join(b.text for b in blocks)
        return blocks, page_count, self._detect_language(combined)

    def _extract_pdf_page_blocks(self, page_dict: dict[str, Any], page_number: int) -> list[DocumentBlock]:
        """Extract blocks from PyMuPDF dict output using font size heuristics."""
        blocks: list[DocumentBlock] = []
        font_sizes: list[float] = []

        for block in page_dict.get("blocks", []):
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    size = span.get("size", 0)
                    if size:
                        font_sizes.append(size)

        median_size = sorted(font_sizes)[len(font_sizes) // 2] if font_sizes else 12.0

        for block in page_dict.get("blocks", []):
            if block.get("type") != 0:
                continue
            lines_text: list[str] = []
            max_size = 0.0
            for line in block.get("lines", []):
                line_parts = []
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if text:
                        line_parts.append(text)
                        max_size = max(max_size, span.get("size", 0))
                if line_parts:
                    lines_text.append(" ".join(line_parts))

            block_text = "\n".join(lines_text).strip()
            if not block_text or self._is_noise_line(block_text):
                continue

            if max_size >= median_size * 1.3 and len(block_text) < 200:
                level = 1 if max_size >= median_size * 1.6 else 2
                block_type = BlockType.HEADING if level == 1 else BlockType.SUBHEADING
                blocks.append(DocumentBlock(
                    block_type=block_type, text=block_text, level=level, page_number=page_number
                ))
            elif self._looks_like_table(block_text):
                blocks.append(DocumentBlock(
                    block_type=BlockType.TABLE, text=block_text, page_number=page_number
                ))
            elif self._looks_like_list(block_text):
                blocks.append(DocumentBlock(
                    block_type=BlockType.LIST, text=block_text, page_number=page_number
                ))
            else:
                blocks.append(DocumentBlock(
                    block_type=BlockType.PARAGRAPH, text=block_text, page_number=page_number
                ))

        return blocks

    def _parse_docx(
        self, content: bytes, filename: str
    ) -> tuple[list[DocumentBlock], int, str]:
        from docx import Document as DocxDocument

        try:
            document = DocxDocument(io.BytesIO(content))
        except Exception as exc:
            raise CorruptedFileError(f"DOCX {filename!r} is corrupted: {exc}") from exc

        blocks: list[DocumentBlock] = []
        heading_style_names = {"Title", "Heading 1", "Heading 2", "Heading 3",
                               "Heading 4", "Heading 5", "Heading 6"}

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name if paragraph.style else ""
            if style_name in heading_style_names or style_name.startswith("Heading"):
                level = self._heading_level_from_style(style_name)
                block_type = BlockType.HEADING if level <= 1 else BlockType.SUBHEADING
                blocks.append(DocumentBlock(
                    block_type=block_type, text=text, level=level
                ))
            elif self._looks_like_list(text):
                blocks.append(DocumentBlock(block_type=BlockType.LIST, text=text))
            else:
                blocks.append(DocumentBlock(block_type=BlockType.PARAGRAPH, text=text))

        for table in document.tables:
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    deduped = self._dedupe_adjacent(cells)
                    rows.append(" | ".join(deduped))
            if rows:
                blocks.append(DocumentBlock(
                    block_type=BlockType.TABLE,
                    text="\n".join(rows),
                    metadata={"row_count": len(rows)},
                ))

        if not blocks:
            raise CorruptedFileError(f"DOCX {filename!r} has no extractable text.")

        combined = "\n\n".join(b.text for b in blocks)
        return blocks, 0, self._detect_language(combined)

    def _parse_xlsx(
        self, content: bytes, filename: str
    ) -> tuple[list[DocumentBlock], int, str]:
        import openpyxl

        try:
            workbook = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)
        except Exception as exc:
            raise CorruptedFileError(f"XLSX {filename!r} is corrupted: {exc}") from exc

        blocks: list[DocumentBlock] = []
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                table_text = f"Sheet: {sheet_name}\n" + "\n".join(rows)
                blocks.append(DocumentBlock(
                    block_type=BlockType.TABLE,
                    text=table_text,
                    metadata={"sheet": sheet_name, "row_count": len(rows)},
                ))

        workbook.close()
        if not blocks:
            raise CorruptedFileError(f"XLSX {filename!r} has no extractable data.")

        combined = "\n\n".join(b.text for b in blocks)
        return blocks, 0, self._detect_language(combined)

    def _parse_pptx(
        self, content: bytes, filename: str
    ) -> tuple[list[DocumentBlock], int, str]:
        from pptx import Presentation

        try:
            prs = Presentation(io.BytesIO(content))
        except Exception as exc:
            raise CorruptedFileError(f"PPTX {filename!r} is corrupted: {exc}") from exc

        blocks: list[DocumentBlock] = []
        slide_count = len(prs.slides)

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_title = ""
            slide_body: list[str] = []

            for shape in slide.shapes:
                if not hasattr(shape, "text") or not shape.text.strip():
                    continue
                text = shape.text.strip()
                if not slide_title and shape == slide.shapes[0]:
                    slide_title = text
                    blocks.append(DocumentBlock(
                        block_type=BlockType.HEADING,
                        text=text,
                        level=1,
                        page_number=slide_idx,
                    ))
                elif self._looks_like_list(text):
                    blocks.append(DocumentBlock(
                        block_type=BlockType.LIST, text=text, page_number=slide_idx
                    ))
                else:
                    slide_body.append(text)

            if slide_body:
                blocks.append(DocumentBlock(
                    block_type=BlockType.PARAGRAPH,
                    text="\n".join(slide_body),
                    page_number=slide_idx,
                ))

        if not blocks:
            raise CorruptedFileError(f"PPTX {filename!r} has no extractable text.")

        combined = "\n\n".join(b.text for b in blocks)
        return blocks, slide_count, self._detect_language(combined)

    def _parse_html(
        self, content: bytes, filename: str
    ) -> tuple[list[DocumentBlock], int, str]:
        text = self._decode_text(content, filename)

        class _HTMLExtractor(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self.blocks: list[DocumentBlock] = []
                self._current_tag = ""
                self._buffer: list[str] = []
                self._in_script = False
                self._in_style = False
                self._in_nav = False

            def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
                self._flush_buffer()
                self._current_tag = tag
                if tag in ("script", "style"):
                    self._in_script = tag == "script"
                    self._in_style = tag == "style"
                if tag == "nav":
                    self._in_nav = True

            def handle_endtag(self, tag: str) -> None:
                self._flush_buffer()
                if tag in ("script", "style"):
                    self._in_script = False
                    self._in_style = False
                if tag == "nav":
                    self._in_nav = False
                self._current_tag = ""

            def handle_data(self, data: str) -> None:
                if self._in_script or self._in_style or self._in_nav:
                    return
                stripped = data.strip()
                if stripped:
                    self._buffer.append(stripped)

            def _flush_buffer(self) -> None:
                if not self._buffer:
                    return
                text = " ".join(self._buffer).strip()
                self._buffer = []
                if not text or _is_navigation_noise(text):
                    return
                tag = self._current_tag
                if tag in ("h1",):
                    self.blocks.append(DocumentBlock(block_type=BlockType.HEADING, text=text, level=1))
                elif tag in ("h2", "h3", "h4", "h5", "h6"):
                    level = int(tag[1])
                    self.blocks.append(DocumentBlock(
                        block_type=BlockType.SUBHEADING, text=text, level=level
                    ))
                elif tag in ("li",):
                    self.blocks.append(DocumentBlock(block_type=BlockType.LIST, text=text))
                elif tag in ("pre", "code"):
                    self.blocks.append(DocumentBlock(block_type=BlockType.CODE, text=text))
                elif tag == "table":
                    self.blocks.append(DocumentBlock(block_type=BlockType.TABLE, text=text))
                else:
                    self.blocks.append(DocumentBlock(block_type=BlockType.PARAGRAPH, text=text))

        parser = _HTMLExtractor()
        parser.feed(text)
        parser._flush_buffer()

        if not parser.blocks:
            # Fallback: strip tags and parse as plain text.
            stripped = re.sub(r"<[^>]+>", "\n", text)
            blocks = self._parse_plain_text_blocks(stripped)
        else:
            blocks = parser.blocks

        combined = "\n\n".join(b.text for b in blocks)
        return blocks, 0, self._detect_language(combined)

    def _parse_markdown_blocks(self, text: str) -> list[DocumentBlock]:
        """Parse markdown into structured blocks."""
        blocks: list[DocumentBlock] = []

        # Extract FAQ pairs first (line-oriented Q/A format).
        faq_blocks = self._extract_faq_blocks(text)
        if faq_blocks:
            return faq_blocks

        lines = text.split("\n")
        line_idx = 0
        while line_idx < len(lines):
            line = lines[line_idx]

            # Heading (markdown).
            heading_match = _MD_HEADING_RE.match(line)
            if heading_match:
                level = len(heading_match.group(1))
                heading_text = heading_match.group(2).strip()
                block_type = BlockType.HEADING if level <= 1 else BlockType.SUBHEADING
                blocks.append(DocumentBlock(block_type=block_type, text=heading_text, level=level))
                line_idx += 1
                continue

            # Table block (consecutive pipe rows).
            if _MD_TABLE_ROW_RE.match(line):
                table_lines = [line]
                line_idx += 1
                while line_idx < len(lines) and _MD_TABLE_ROW_RE.match(lines[line_idx]):
                    table_lines.append(lines[line_idx])
                    line_idx += 1
                blocks.append(DocumentBlock(
                    block_type=BlockType.TABLE, text="\n".join(table_lines)
                ))
                continue

            # List block (consecutive list items).
            if _MD_LIST_RE.match(line):
                list_lines = [line]
                line_idx += 1
                while line_idx < len(lines) and _MD_LIST_RE.match(lines[line_idx]):
                    list_lines.append(lines[line_idx])
                    line_idx += 1
                blocks.append(DocumentBlock(
                    block_type=BlockType.LIST, text="\n".join(list_lines)
                ))
                continue

            # Fenced code block.
            if line.strip().startswith("```"):
                fence_lines = [line]
                line_idx += 1
                while line_idx < len(lines):
                    fence_lines.append(lines[line_idx])
                    if lines[line_idx].strip().startswith("```") and len(fence_lines) > 1:
                        line_idx += 1
                        break
                    line_idx += 1
                code_text = "\n".join(fence_lines)
                inner = _MD_CODE_FENCE_RE.search(code_text)
                blocks.append(DocumentBlock(
                    block_type=BlockType.CODE,
                    text=inner.group(1).strip() if inner else code_text,
                ))
                continue

            # Blank line — skip.
            if not line.strip():
                line_idx += 1
                continue

            # Paragraph: accumulate until structural boundary.
            para_lines = [line]
            line_idx += 1
            while line_idx < len(lines):
                next_line = lines[line_idx]
                if (
                    not next_line.strip()
                    or _MD_HEADING_RE.match(next_line)
                    or _MD_TABLE_ROW_RE.match(next_line)
                    or _MD_LIST_RE.match(next_line)
                    or next_line.strip().startswith("```")
                ):
                    break
                para_lines.append(next_line)
                line_idx += 1

            para_text = "\n".join(para_lines).strip()
            if para_text and not self._is_noise_line(para_text):
                if self._looks_like_code(para_text):
                    blocks.extend(self._split_code_blocks(para_text))
                else:
                    blocks.append(DocumentBlock(block_type=BlockType.PARAGRAPH, text=para_text))

        return self._merge_consecutive_lists(blocks)

    def _parse_plain_text_blocks(self, text: str) -> list[DocumentBlock]:
        blocks: list[DocumentBlock] = []
        for para in self._split_paragraphs(text):
            if self._is_noise_line(para):
                continue
            if self._looks_like_faq(para):
                blocks.extend(self._extract_faq_blocks(para))
            elif self._looks_like_table(para):
                blocks.append(DocumentBlock(block_type=BlockType.TABLE, text=para))
            elif self._looks_like_list(para):
                blocks.append(DocumentBlock(block_type=BlockType.LIST, text=para))
            elif self._looks_like_code(para):
                blocks.extend(self._split_code_blocks(para))
            else:
                blocks.append(DocumentBlock(block_type=BlockType.PARAGRAPH, text=para))
        return blocks

    def _extract_faq_blocks(self, text: str) -> list[DocumentBlock]:
        """Extract Q+A pairs as individual FAQ blocks."""
        blocks: list[DocumentBlock] = []
        lines = text.splitlines()
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            q_match = _FAQ_Q_RE.match(line)
            if q_match:
                question = q_match.group(1).strip()
                answer_lines: list[str] = []
                i += 1
                while i < len(lines):
                    next_line = lines[i].strip()
                    if _FAQ_Q_RE.match(next_line):
                        break
                    a_match = _FAQ_A_RE.match(next_line)
                    if a_match:
                        answer_lines.append(a_match.group(1).strip())
                    elif answer_lines:
                        answer_lines.append(next_line)
                    i += 1
                answer = " ".join(answer_lines).strip()
                if question and answer:
                    blocks.append(DocumentBlock(
                        block_type=BlockType.FAQ,
                        text=f"Q: {question}\nA: {answer}",
                        metadata={"question": question, "answer": answer},
                    ))
                continue
            i += 1
        return blocks

    def _split_code_blocks(self, text: str) -> list[DocumentBlock]:
        """Split source code into one function/class per block."""
        blocks: list[DocumentBlock] = []
        boundaries: list[tuple[int, str]] = []

        for match in _PY_DEF_RE.finditer(text):
            boundaries.append((match.start(), "function"))
        for match in _PY_CLASS_RE.finditer(text):
            boundaries.append((match.start(), "class"))

        if not boundaries:
            return [DocumentBlock(block_type=BlockType.CODE, text=text)]

        boundaries.sort(key=lambda x: x[0])
        for idx, (start, kind) in enumerate(boundaries):
            end = boundaries[idx + 1][0] if idx + 1 < len(boundaries) else len(text)
            code_text = text[start:end].strip()
            if code_text:
                blocks.append(DocumentBlock(
                    block_type=BlockType.CODE, text=code_text, metadata={"kind": kind}
                ))
        return blocks

    def _remove_noise(self, blocks: list[DocumentBlock]) -> list[DocumentBlock]:
        """Filter headers, footers, page numbers, watermarks, blank content."""
        cleaned: list[DocumentBlock] = []
        for block in blocks:
            text = block.text.strip()
            if not text:
                continue
            if block.block_type in (BlockType.TABLE, BlockType.CODE, BlockType.FAQ, BlockType.LIST):
                cleaned.append(block)
                continue
            # Remove page-number-only and watermark lines.
            lines = [ln for ln in text.splitlines() if not self._is_noise_line(ln)]
            filtered = "\n".join(lines).strip()
            if filtered and len(filtered) >= 3:
                cleaned.append(block.model_copy(update={"text": filtered}))
        return cleaned

    @staticmethod
    def _is_noise_line(line: str) -> bool:
        stripped = line.strip()
        if not stripped:
            return True
        if _PAGE_NUMBER_RE.match(stripped):
            return True
        if _HEADER_FOOTER_RE.match(stripped):
            return True
        if _WATERMARK_RE.search(stripped) and len(stripped) < 40:
            return True
        return False

    @staticmethod
    def _looks_like_table(text: str) -> bool:
        lines = text.splitlines()
        pipe_lines = sum(1 for ln in lines if "|" in ln and ln.count("|") >= 2)
        tab_lines = sum(1 for ln in lines if "\t" in ln)
        return pipe_lines >= 2 or tab_lines >= 2

    @staticmethod
    def _looks_like_list(text: str) -> bool:
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        if len(lines) < 2:
            return bool(_MD_LIST_RE.match(text))
        list_lines = sum(
            1 for ln in lines
            if re.match(r"^[-*+•]\s+", ln) or re.match(r"^\d+[.)]\s+", ln)
        )
        return list_lines >= len(lines) * 0.6

    @staticmethod
    def _looks_like_code(text: str) -> bool:
        indicators = ("def ", "class ", "import ", "function ", "const ", "var ", "```")
        return any(ind in text for ind in indicators)

    @staticmethod
    def _looks_like_faq(text: str) -> bool:
        return bool(_FAQ_Q_RE.search(text))

    @staticmethod
    def _split_paragraphs(text: str) -> list[str]:
        return [p.strip() for p in re.split(r"\n\s*\n+", text) if p.strip()]

    @staticmethod
    def _find_next_structure(text: str, pos: int) -> int:
        patterns = [_MD_HEADING_RE, _MD_TABLE_ROW_RE, _MD_LIST_RE, _MD_CODE_FENCE_RE]
        earliest = len(text)
        for pattern in patterns:
            match = pattern.search(text, pos + 1)
            if match and match.start() < earliest:
                earliest = match.start()
        return earliest

    @staticmethod
    def _heading_level_from_style(style_name: str) -> int:
        if style_name == "Title":
            return 1
        if style_name.startswith("Heading"):
            try:
                return int(style_name.split()[-1])
            except ValueError:
                return 2
        return 2

    @staticmethod
    def _dedupe_adjacent(items: list[str]) -> list[str]:
        result: list[str] = []
        for item in items:
            if not result or result[-1] != item:
                result.append(item)
        return result

    @staticmethod
    def _merge_consecutive_lists(blocks: list[DocumentBlock]) -> list[DocumentBlock]:
        if not blocks:
            return blocks
        merged: list[DocumentBlock] = [blocks[0]]
        for block in blocks[1:]:
            prev = merged[-1]
            if prev.block_type == BlockType.LIST and block.block_type == BlockType.LIST:
                merged[-1] = prev.model_copy(update={"text": prev.text + "\n" + block.text})
            else:
                merged.append(block)
        return merged

    @staticmethod
    def _decode_text(content: bytes, filename: str) -> str:
        for encoding in ("utf-8", "utf-8-sig", "latin-1"):
            try:
                text = content.decode(encoding)
                if text.strip():
                    return text
            except UnicodeDecodeError:
                continue
        raise CorruptedFileError(f"Text file {filename!r} uses unsupported encoding.")

    @staticmethod
    def _detect_language(text: str) -> str:
        """Simple language detection without external models."""
        if not text:
            return "en"
        # CJK character ranges.
        cjk = sum(1 for c in text if "\u4e00" <= c <= "\u9fff" or "\u3040" <= c <= "\u30ff")
        if cjk > len(text) * 0.1:
            return "zh"
        # Devanagari.
        devanagari = sum(1 for c in text if "\u0900" <= c <= "\u097f")
        if devanagari > len(text) * 0.1:
            return "hi"
        try:
            from langdetect import detect
            return detect(text[:2000])
        except Exception:
            return "en"


def _is_navigation_noise(text: str) -> bool:
    lower = text.lower()
    nav_terms = ("home", "menu", "login", "sign up", "cookie", "privacy policy")
    return any(term in lower for term in nav_terms) and len(text) < 60


# Module-level convenience for dependency injection.
_default_parser: DocumentParser | None = None


def get_document_parser() -> DocumentParser:
    global _default_parser
    if _default_parser is None:
        _default_parser = DocumentParser()
    return _default_parser
