import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Paths
ASSETS_DIR = r"d:\local-rag\docs\presentation_assets"
OUTPUT_PATH = r"d:\local-rag\docs\SipraHub_Local_RAG_v3.pptx"

# SipraHub Brand Palette from sipra-2026-sales-deck 2.pptx
BRAND_RED = RGBColor(206, 33, 36)        # #CE2124
COLOR_TITLE = RGBColor(26, 26, 26)       # #1A1A1A
COLOR_BODY = RGBColor(61, 61, 61)        # #3D3D3D
COLOR_MUTED = RGBColor(122, 122, 122)    # #7A7A7A
COLOR_BORDER = RGBColor(229, 231, 235)   # #E5E7EB
COLOR_CARD_BG = RGBColor(248, 249, 250)  # #F8F9FA
COLOR_WHITE = RGBColor(255, 255, 255)    # #FFFFFF
COLOR_DARK = RGBColor(15, 23, 42)        # #0F172A
COLOR_LIGHT_RED = RGBColor(254, 242, 242)# #FEF2F2

FONT_NAME = "Calibri"

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, breadcrumb="CASE STUDY 06.", category="100% FULLY LOCAL RAG"):
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.22), Inches(6.0), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = breadcrumb + "  "
        r1.font.name = FONT_NAME
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        r2 = p.add_run()
        r2.text = f"|  {category}"
        r2.font.name = FONT_NAME
        r2.font.size = Pt(9.5)
        r2.font.bold = True
        r2.font.color.rgb = COLOR_MUTED

    def add_footer(slide, slide_num):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.22), Inches(8.8), Inches(0.015))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BORDER
        line.line.color.rgb = COLOR_BORDER

        tb_left = slide.shapes.add_textbox(Inches(0.6), Inches(5.26), Inches(2.5), Inches(0.25))
        tf_l = tb_left.text_frame
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
        p_l = tf_l.paragraphs[0]
        r_l = p_l.add_run()
        r_l.text = "SipraHub"
        r_l.font.name = FONT_NAME
        r_l.font.size = Pt(10)
        r_l.font.bold = True
        r_l.font.color.rgb = BRAND_RED

        tb_center = slide.shapes.add_textbox(Inches(3.2), Inches(5.26), Inches(3.6), Inches(0.25))
        tf_c = tb_center.text_frame
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.CENTER
        r_c = p_c.add_run()
        r_c.text = "AI-Driven Digital Transformation"
        r_c.font.name = FONT_NAME
        r_c.font.size = Pt(9.5)
        r_c.font.color.rgb = COLOR_MUTED

        tb_right = slide.shapes.add_textbox(Inches(8.8), Inches(5.26), Inches(0.6), Inches(0.25))
        tf_r = tb_right.text_frame
        tf_r.margin_left = tf_r.margin_top = tf_r.margin_right = tf_r.margin_bottom = 0
        p_r = tf_r.paragraphs[0]
        p_r.alignment = PP_ALIGN.RIGHT
        r_r = p_r.add_run()
        r_r.text = str(slide_num)
        r_r.font.name = FONT_NAME
        r_r.font.size = Pt(10)
        r_r.font.bold = True
        r_r.font.color.rgb = COLOR_TITLE

    # =========================================================================
    # SLIDE 1 — COVER SLIDE (100% FULLY LOCAL RAG)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    
    tb1 = s1.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(4.5), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

    p = tf1.paragraphs[0]
    r = p.add_run()
    r.text = "ENTERPRISE AI PORTFOLIO"
    r.font.name = FONT_NAME
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = BRAND_RED

    p2 = tf1.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = "SipraHub Fully Local RAG"
    r2.font.name = FONT_NAME
    r2.font.size = Pt(28)
    r2.font.bold = True
    r2.font.color.rgb = COLOR_TITLE

    p3 = tf1.add_paragraph()
    p3.space_before = Pt(4)
    r3 = p3.add_run()
    r3.text = "100% On-Premise Sovereign Knowledge Retrieval & Air-Gapped AI"
    r3.font.name = FONT_NAME
    r3.font.size = Pt(13)
    r3.font.color.rgb = COLOR_BODY

    p4 = tf1.add_paragraph()
    p4.space_before = Pt(16)
    r4 = p4.add_run()
    r4.text = "A completely air-gapped enterprise RAG platform running 100% locally on private infrastructure with zero cloud dependencies, zero external API costs, and zero data leakage."
    r4.font.name = FONT_NAME
    r4.font.size = Pt(11)
    r4.font.color.rgb = COLOR_MUTED

    p5 = tf1.add_paragraph()
    p5.space_before = Pt(22)
    r5 = p5.add_run()
    r5.text = "CATEGORY: "
    r5.font.name = FONT_NAME
    r5.font.size = Pt(9.5)
    r5.font.bold = True
    r5.font.color.rgb = COLOR_TITLE
    r5_b = p5.add_run()
    r5_b.text = "Fully Local RAG | 100% On-Premise | Air-Gapped Data Sovereignty"
    r5_b.font.name = FONT_NAME
    r5_b.font.size = Pt(9.5)
    r5_b.font.color.rgb = BRAND_RED

    hero_path = os.path.join(ASSETS_DIR, "hero_app.png")
    if os.path.exists(hero_path):
        s1.shapes.add_picture(hero_path, Inches(5.2), Inches(0.7), Inches(4.3), Inches(4.2))

    add_footer(s1, 1)

    # =========================================================================
    # SLIDE 2 — EXECUTIVE SUMMARY
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "CASE STUDY 06.", "EXECUTIVE SUMMARY")
    
    tb = s2.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Executive Summary: 100% Sovereign On-Premise Intelligence"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    p_sub = tf.add_paragraph()
    p_sub.space_before = Pt(4)
    r_sub = p_sub.add_run()
    r_sub.text = "SipraHub developed a 100% Fully Local Retrieval-Augmented Generation platform. Every component—from document ingestion and vector embeddings to neural reranking and LLM inference—executes strictly within private on-premise hardware."
    r_sub.font.name = FONT_NAME
    r_sub.font.size = Pt(11)
    r_sub.font.color.rgb = COLOR_BODY

    metrics = [
        ("100%", "Air-Gapped Privacy", "Zero external network requests; zero document text or embeddings leave the local LAN."),
        ("$0", "Recurring API Cost", "Eliminates monthly third-party cloud LLM subscription and token billing completely."),
        ("0", "Fabricated Policies", "Strict local anti-hallucination prompt guardrails suppress unstated rules."),
        ("13+", "Document Formats", "Local parsing for PDF, DOCX, XLSX, PPTX, Markdown, and text on private disks.")
    ]
    x_card = 0.6
    for val, label, subtext in metrics:
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_card), Inches(2.05), Inches(2.05), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER

        tb_c = s2.shapes.add_textbox(Inches(x_card + 0.15), Inches(2.2), Inches(1.75), Inches(2.4))
        tfc = tb_c.text_frame
        tfc.word_wrap = True
        
        p1 = tfc.paragraphs[0]
        r1 = p1.add_run()
        r1.text = val
        r1.font.name = FONT_NAME
        r1.font.size = Pt(32)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tfc.add_paragraph()
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = FONT_NAME
        r2.font.size = Pt(12)
        r2.font.bold = True
        r2.font.color.rgb = COLOR_TITLE

        p3 = tfc.add_paragraph()
        p3.space_before = Pt(6)
        r3 = p3.add_run()
        r3.text = subtext
        r3.font.name = FONT_NAME
        r3.font.size = Pt(10)
        r3.font.color.rgb = COLOR_MUTED

        x_card += 2.25

    add_footer(s2, 2)

    # =========================================================================
    # SLIDE 3 — BUSINESS CHALLENGE (WHY CLOUD AI FAILS ENTERPRISES)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "CASE STUDY 06.", "BUSINESS CHALLENGE")

    tb = s3.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Why Public Cloud AI Violates Enterprise Sovereignty"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    p_sub = tf.add_paragraph()
    r_sub = p_sub.add_run()
    r_sub.text = "Sending internal business documents, PRDs, and personnel frameworks to cloud LLMs creates immense organizational liability."
    r_sub.font.name = FONT_NAME
    r_sub.font.size = Pt(11)
    r_sub.font.color.rgb = COLOR_BODY

    challenges = [
        ("01", "Confidential IP Exposure", "Uploading proprietary source code, roadmap documents, and HR policies to cloud AI endpoints violates corporate governance and client confidentiality agreements."),
        ("02", "Regulatory Compliance Non-Viability", "Defense, finance, healthcare, and enterprise IT contracts strictly mandate air-gapped data residency that public cloud AI APIs cannot satisfy."),
        ("03", "Runaway Token Billing", "Enterprise-wide document queries generate millions of tokens per month, resulting in unpredictable and exorbitant third-party API subscription costs."),
        ("04", "External Outages & Cloud Dependency", "Cloud AI services experience frequent latency degradation and downtime, stalling critical internal employee workflows and operations.")
    ]
    x_pos = 0.6
    for num, title, desc in challenges:
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.9), Inches(2.05), Inches(3.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER

        tb_c = s3.shapes.add_textbox(Inches(x_pos + 0.15), Inches(2.05), Inches(1.75), Inches(2.6))
        tfc = tb_c.text_frame
        tfc.word_wrap = True

        p1 = tfc.paragraphs[0]
        r1 = p1.add_run()
        r1.text = num
        r1.font.name = FONT_NAME
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tfc.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = title
        r2.font.name = FONT_NAME
        r2.font.size = Pt(12)
        r2.font.bold = True
        r2.font.color.rgb = COLOR_TITLE

        p3 = tfc.add_paragraph()
        p3.space_before = Pt(6)
        r3 = p3.add_run()
        r3.text = desc
        r3.font.name = FONT_NAME
        r3.font.size = Pt(10)
        r3.font.color.rgb = COLOR_BODY

        x_pos += 2.25

    add_footer(s3, 3)

    # =========================================================================
    # SLIDE 4 — OUR SOLUTION (100% LOCAL STACK)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "CASE STUDY 06.", "OUR SOLUTION")

    tb = s4.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(4.5), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "SipraHub 100% Local RAG Solution"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = "A complete, self-contained generative retrieval engine designed to run entirely on private company infrastructure with zero cloud exposure."
    r2.font.name = FONT_NAME
    r2.font.size = Pt(11)
    r2.font.color.rgb = COLOR_BODY

    pillars = [
        ("OmniRoute Local AI Gateway", "Routes all LLM inference through a local gateway (http://127.0.0.1:20128/v1) running Qwen 3: 8B — zero external token billing."),
        ("nomic-embed-text 768d Embeddings", "Generates dense 768-dimensional semantic vectors locally via Ollama; written directly to pgvector — never leaves the private network."),
        ("On-Premise PostgreSQL + pgvector", "Stores all document chunks and VECTOR(768) embeddings inside a private on-premise PostgreSQL instance with tsvector full-text indexing."),
        ("FlashRank Neural Cross-Encoder", "ms-marco-TinyBERT-L-2-v2 reranks top-30 hybrid candidates to top-10 high-precision passages fully on CPU.")
    ]
    for p_title, p_desc in pillars:
        p_item = tf.add_paragraph()
        p_item.space_before = Pt(10)
        r_t = p_item.add_run()
        r_t.text = f"{p_title}: "
        r_t.font.name = FONT_NAME
        r_t.font.size = Pt(10.5)
        r_t.font.bold = True
        r_t.font.color.rgb = BRAND_RED

        r_d = p_item.add_run()
        r_d.text = p_desc
        r_d.font.name = FONT_NAME
        r_d.font.size = Pt(10)
        r_d.font.color.rgb = COLOR_BODY

    if os.path.exists(hero_path):
        s4.shapes.add_picture(hero_path, Inches(5.3), Inches(0.8), Inches(4.2), Inches(4.1))

    add_footer(s4, 4)

    # =========================================================================
    # SLIDE 5 — PRODUCT EXPERIENCE (LOCAL INTERFACE)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "CASE STUDY 06.", "PRODUCT EXPERIENCE")

    tb = s5.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Enterprise Workspace Running on Local Intranet"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    doc_path = os.path.join(ASSETS_DIR, "doc_management.png")
    if os.path.exists(doc_path):
        s5.shapes.add_picture(doc_path, Inches(0.6), Inches(1.4), Inches(5.2), Inches(3.5))

    callouts = [
        ("1. Local Document Vault", "Employees upload documents stored directly on local private disks, isolated from the public internet."),
        ("2. Real-Time Local Streaming", "Responses stream word-by-word with immediate feedback directly from the local Ollama process."),
        ("3. Verifiable Local Citations", "Every response features clickable source pills mapping directly to local file chunks and version IDs."),
        ("4. Air-Gapped Anti-Hallucination", "Strict grounding prompts enforce that unstated company policies are refused rather than fabricated.")
    ]
    y_call = 1.4
    for c_title, c_desc in callouts:
        c_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(y_call), Inches(3.4), Inches(0.78))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = COLOR_CARD_BG
        c_box.line.color.rgb = COLOR_BORDER

        tb_call = s5.shapes.add_textbox(Inches(6.12), Inches(y_call + 0.08), Inches(3.15), Inches(0.62))
        tf_c = tb_call.text_frame
        tf_c.word_wrap = True

        p1 = tf_c.paragraphs[0]
        r1 = p1.add_run()
        r1.text = c_title
        r1.font.name = FONT_NAME
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tf_c.add_paragraph()
        p2.space_before = Pt(2)
        r2 = p2.add_run()
        r2.text = c_desc
        r2.font.name = FONT_NAME
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_BODY

        y_call += 0.88

    add_footer(s5, 5)

    # =========================================================================
    # SLIDE 6 — KEY FEATURES (FULLY LOCAL CAPABILITIES)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "CASE STUDY 06.", "LOCAL RAG CAPABILITIES")

    tb = s6.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "100% Fully Local Enterprise Capabilities"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    features = [
        ("Agentic 10-Route Intent Router", "Rule-based router (zero LLM overhead) classifies queries into 10 routes: DOCUMENT_QA, SUMMARY, WEB, CALCULATOR, HYBRID, GENERAL, DIRECT, and more."),
        ("Hybrid RRF Search + Reranking", "Reciprocal Rank Fusion merges pgvector cosine + tsvector keyword results; FlashRank TinyBERT reranks top-30 → top-10 on CPU."),
        ("OmniRoute Streaming Generation", "Real-time token-by-token SSE streaming via local OmniRoute gateway; thinking-block filtered before client delivery."),
        ("13-Format Local Document Parsing", "Native in-process parsing: PDF (PyMuPDF+pypdf), DOCX, XLSX, PPTX, MD, CSV, JSON, LOG, HTML — zero cloud converters."),
        ("Verifiable Source Citations", "Every answer embeds chunk_id, document_version_id, similarity score, and section title for audit-trail compliance."),
        ("Long-Term Conversation Memory", "Rule-based extractor captures user facts and preferences post-response; top-5 memories injected via cosine retrieval (threshold 0.75).")
    ]
    coords = [
        (0.6, 1.5), (3.6, 1.5), (6.6, 1.5),
        (0.6, 3.3), (3.6, 3.3), (6.6, 3.3)
    ]
    for i, (title, desc) in enumerate(features):
        x, y = coords[i]
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER

        tb_f = s6.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12), Inches(2.5), Inches(1.35))
        tf_f = tb_f.text_frame
        tf_f.word_wrap = True

        p1 = tf_f.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = FONT_NAME
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tf_f.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = FONT_NAME
        r2.font.size = Pt(10)
        r2.font.color.rgb = COLOR_BODY

    add_footer(s6, 6)

    # =========================================================================
    # SLIDE 7 — TECHNICAL ARCHITECTURE (AIR-GAPPED)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "CASE STUDY 06.", "AIR-GAPPED ARCHITECTURE")

    tb = s7.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "100% On-Premise Air-Gapped Architecture"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    arch_path = os.path.join(ASSETS_DIR, "architecture_diagram.png")
    if os.path.exists(arch_path):
        s7.shapes.add_picture(arch_path, Inches(0.6), Inches(1.35), Inches(8.8), Inches(3.65))

    add_footer(s7, 7)

    # =========================================================================
    # SLIDE 8 — WORKFLOW (LOCAL PIPELINE)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "CASE STUDY 06.", "LOCAL WORKFLOW")

    tb = s8.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Local Ingestion & Query Resolution Lifecycle"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    steps = [
        ("1. Upload & Store", "File stored on private local disk; SHA-256 content hash computed locally to prevent duplicate re-embedding."),
        ("2. Parse & Chunk", "PyMuPDF/pypdf/docx extraction → 1500-char sliding window (300 overlap) with paragraph-boundary preference."),
        ("3. Embed & Index", "nomic-embed-text 768d vectors generated locally by Ollama; stored in on-premise pgvector VECTOR(768) column."),
        ("4. Intent Route", "Rule-based router classifies query into 1 of 10 routes in <1ms with zero LLM overhead; DOCUMENT_QA triggers hybrid retrieval."),
        ("5. Hybrid RRF Rerank", "pgvector cosine + tsvector merged via RRF; FlashRank TinyBERT reranks top-30 → top-10 on CPU."),
        ("6. Stream Answer", "OmniRoute Qwen 3: 8B streams tokens via SSE; thinking-tags filtered real-time; citations appended.")
    ]
    x_s = 0.6
    for i, (title, desc) in enumerate(steps):
        s_card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_s), Inches(1.7), Inches(1.35), Inches(3.2))
        s_card.fill.solid()
        s_card.fill.fore_color.rgb = COLOR_CARD_BG
        s_card.line.color.rgb = BRAND_RED if i in (3, 4, 5) else COLOR_BORDER

        tb_s = s8.shapes.add_textbox(Inches(x_s + 0.08), Inches(1.8), Inches(1.2), Inches(3.0))
        tf_s = tb_s.text_frame
        tf_s.word_wrap = True

        p1 = tf_s.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = FONT_NAME
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tf_s.add_paragraph()
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = FONT_NAME
        r2.font.size = Pt(9)
        r2.font.color.rgb = COLOR_BODY

        x_s += 1.48

    add_footer(s8, 8)

    # =========================================================================
    # SLIDE 9 — BUSINESS IMPACT (SOVEREIGNTY & COST)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "CASE STUDY 06.", "BUSINESS IMPACT")

    tb = s9.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Tangible Business Value of Fully Local Architecture"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    impacts = [
        ("$0", "Monthly API Billing", "Zero recurring token subscription fees, shielding enterprise budgets from fluctuating cloud AI pricing."),
        ("100%", "Data Sovereignty", "Guarantees complete compliance with strict internal IP policies, defense standards, and privacy mandates."),
        ("0", "Fabricated Policies", "Strict local document grounding eliminates employee confusion caused by hallucinated HR rules."),
        ("0 B", "External Data Egress", "Zero bytes leave the private enterprise network, guaranteeing total confidentiality.")
    ]
    x_i = 0.6
    for metric, title, desc in impacts:
        card = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_i), Inches(1.8), Inches(2.05), Inches(3.1))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER

        tb_i = s9.shapes.add_textbox(Inches(x_i + 0.15), Inches(1.95), Inches(1.75), Inches(2.8))
        tf_i = tb_i.text_frame
        tf_i.word_wrap = True

        p1 = tf_i.paragraphs[0]
        r1 = p1.add_run()
        r1.text = metric
        r1.font.name = FONT_NAME
        r1.font.size = Pt(30)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tf_i.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = title
        r2.font.name = FONT_NAME
        r2.font.size = Pt(12)
        r2.font.bold = True
        r2.font.color.rgb = COLOR_TITLE

        p3 = tf_i.add_paragraph()
        p3.space_before = Pt(6)
        r3 = p3.add_run()
        r3.text = desc
        r3.font.name = FONT_NAME
        r3.font.size = Pt(10)
        r3.font.color.rgb = COLOR_BODY

        x_i += 2.25

    add_footer(s9, 9)

    # =========================================================================
    # SLIDE 10 — SECURITY & AIR-GAPPED GOVERNANCE
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "CASE STUDY 06.", "AIR-GAPPED SECURITY")

    tb = s10.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Air-Gapped Security Controls & Enterprise Governance"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    sec_pillars = [
        ("100% Private Vector Store", "Embeddings and document text reside inside private on-premise PostgreSQL tables, completely isolated from multi-tenant cloud AI systems or external training scrapers."),
        ("Cryptographic JWT Authentication", "Secure user sessions authenticated via cryptographically signed JWT tokens with bcrypt password hashing and 32-character minimum secret enforcement."),
        ("Prompt Injection Neutralization", "Local prompt shields instruct the LLM to treat document content strictly as untrusted data, preventing adversarial jailbreaks embedded in user files."),
        ("Zero Telemetry & External Logging", "The entire stack is configured with zero external telemetry, zero ping-backs, and zero data capture, guaranteeing complete organizational privacy.")
    ]
    coords_sec = [(0.6, 1.6), (5.1, 1.6), (0.6, 3.3), (5.1, 3.3)]
    for i, (title, desc) in enumerate(sec_pillars):
        x, y = coords_sec[i]
        card = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER

        tb_s = s10.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(3.9), Inches(1.2))
        tf_s = tb_s.text_frame
        tf_s.word_wrap = True

        p1 = tf_s.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = FONT_NAME
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tf_s.add_paragraph()
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = FONT_NAME
        r2.font.size = Pt(10)
        r2.font.color.rgb = COLOR_BODY

    add_footer(s10, 10)

    # =========================================================================
    # SLIDE 11 — BEFORE VS AFTER (CLOUD DEPENDENCY VS LOCAL RAG)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "CASE STUDY 06.", "BEFORE VS AFTER")

    tb = s11.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Operational Paradigm Shift: Cloud AI vs 100% Local RAG"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    card_before = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(4.3), Inches(3.3))
    card_before.fill.solid()
    card_before.fill.fore_color.rgb = COLOR_CARD_BG
    card_before.line.color.rgb = COLOR_BORDER

    tb_b = s11.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(3.9), Inches(2.9))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    pb = tf_b.paragraphs[0]
    rb = pb.add_run()
    rb.text = "TRADITIONAL / CLOUD AI"
    rb.font.name = FONT_NAME
    rb.font.size = Pt(14)
    rb.font.bold = True
    rb.font.color.rgb = COLOR_MUTED

    b_points = [
        "Confidential documents uploaded to third-party cloud APIs.",
        "Monthly recurring per-token subscription and usage bills.",
        "Public models hallucinating non-existent enterprise rules.",
        "Zero traceability or verified document paragraph citations.",
        "Vulnerable to external vendor service outages and API deprecations."
    ]
    for pt in b_points:
        p_pt = tf_b.add_paragraph()
        p_pt.space_before = Pt(8)
        r_pt = p_pt.add_run()
        r_pt.text = f"[X]  {pt}"
        r_pt.font.name = FONT_NAME
        r_pt.font.size = Pt(10)
        r_pt.font.color.rgb = COLOR_BODY

    card_after = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.6), Inches(4.3), Inches(3.3))
    card_after.fill.solid()
    card_after.fill.fore_color.rgb = COLOR_LIGHT_RED
    card_after.line.color.rgb = BRAND_RED

    tb_a = s11.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(3.9), Inches(2.9))
    tf_a = tb_a.text_frame
    tf_a.word_wrap = True
    pa = tf_a.paragraphs[0]
    ra = pa.add_run()
    ra.text = "SIPRAHUB 100% LOCAL RAG"
    ra.font.name = FONT_NAME
    ra.font.size = Pt(14)
    ra.font.bold = True
    ra.font.color.rgb = BRAND_RED

    a_points = [
        "100% Air-gapped on-premise execution with zero data egress.",
        "Zero token billing; runs on owned private hardware.",
        "Strict anti-hallucination prompts; unstated rules refused.",
        "Deterministic clickable citations with chunk-level auditability.",
        "Autonomous, resilient, and immune to internet or cloud outages."
    ]
    for pt in a_points:
        p_pt = tf_a.add_paragraph()
        p_pt.space_before = Pt(8)
        r_pt = p_pt.add_run()
        r_pt.text = f"[OK]  {pt}"
        r_pt.font.name = FONT_NAME
        r_pt.font.size = Pt(10)
        r_pt.font.bold = True
        r_pt.font.color.rgb = COLOR_TITLE

    add_footer(s11, 11)

    # =========================================================================
    # SLIDE 12 — VALUE DELIVERED (LOCAL WINS)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "CASE STUDY 06.", "VALUE DELIVERED")

    tb = s12.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Enterprise Value Delivered by Fully Local RAG"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    values = [
        ("Absolute Data Sovereignty", "Eliminates legal and compliance barriers by guaranteeing company IP never leaves private servers."),
        ("Predictable TCO & Zero API Fees", "One-time hardware deployment cost with zero recurring SaaS subscriptions or per-query pricing."),
        ("Auditable Fact Grounding", "Every answer cites exact document titles, version hashes, and chunk IDs for verifiable compliance."),
        ("Multi-Format Parsing on Disk", "Direct local ingestion of Word, PDF, Excel, and Markdown files without cloud document converters."),
        ("Paragraph-Aware Chunking", "1500-character sliding window with paragraph preservation ensures complete policy context."),
        ("Self-Contained Local UI", "Responsive React 18 interface with real-time SSE streaming running entirely on the company intranet.")
    ]
    coords_val = [
        (0.6, 1.6), (5.1, 1.6),
        (0.6, 2.8), (5.1, 2.8),
        (0.6, 4.0), (5.1, 4.0)
    ]
    for i, (title, desc) in enumerate(values):
        x, y = coords_val[i]
        card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(1.05))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER

        tb_v = s12.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(4.0), Inches(0.85))
        tf_v = tb_v.text_frame
        tf_v.word_wrap = True

        p1 = tf_v.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = FONT_NAME
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tf_v.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = FONT_NAME
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_BODY

    add_footer(s12, 12)

    # =========================================================================
    # SLIDE 13 — FUTURE ROADMAP (ON-PREMISE SCALING)
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "CASE STUDY 06.", "ON-PREMISE ROADMAP")

    tb = s13.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "On-Premise Enterprise Scaling Roadmap"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    phases = [
        ("PHASE 1 (Q1)", "On-Premise RBAC & Auth", [
            "Document-level Role-Based Access Control (RBAC).",
            "Local Active Directory / LDAP enterprise SSO.",
            "Departmental workspace isolation on private disk."
        ]),
        ("PHASE 2 (Q2)", "Local OCR & Multi-Hop Agents", [
            "PaddleOCR pipeline for scanned physical documents.",
            "Local multi-hop reasoning agents for complex QA.",
            "Local visual analysis for diagrams and blueprints."
        ]),
        ("PHASE 3 (Q3)", "Dedicated GPU Acceleration", [
            "Dedicated on-premise RTX/A-series GPU provisioning.",
            "Sub-10-second end-to-end local generation.",
            "Automated continuous evaluation CI/CD pipeline."
        ])
    ]
    x_ph = 0.6
    for p_num, p_title, items in phases:
        card = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_ph), Inches(1.65), Inches(2.8), Inches(3.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = BRAND_RED if "PHASE 1" in p_num else COLOR_BORDER

        tb_p = s13.shapes.add_textbox(Inches(x_ph + 0.15), Inches(1.8), Inches(2.5), Inches(3.0))
        tf_p = tb_p.text_frame
        tf_p.word_wrap = True

        p1 = tf_p.paragraphs[0]
        r1 = p1.add_run()
        r1.text = p_num
        r1.font.name = FONT_NAME
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tf_p.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = p_title
        r2.font.name = FONT_NAME
        r2.font.size = Pt(13)
        r2.font.bold = True
        r2.font.color.rgb = COLOR_TITLE

        for it in items:
            p_it = tf_p.add_paragraph()
            p_it.space_before = Pt(8)
            r_it = p_it.add_run()
            r_it.text = f"- {it}"
            r_it.font.name = FONT_NAME
            r_it.font.size = Pt(10)
            r_it.font.color.rgb = COLOR_BODY

        x_ph += 3.0

    add_footer(s13, 13)

    # =========================================================================
    # SLIDE 14 — TECHNOLOGY STACK (100% LOCAL TOOLS)
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "CASE STUDY 06.", "LOCAL TECH STACK")

    tb = s14.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(8.8), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "100% Fully Local Technology Stack"
    r.font.name = FONT_NAME
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    stacks = [
        ("Frontend", "React 18 + TypeScript + Vite, Zustand state, Axios, EventSource SSE streaming client"),
        ("Backend API", "FastAPI (Python 3.12), Pydantic v2, SQLAlchemy 2.0 async (psycopg3), Uvicorn ASGI"),
        ("Vector Database", "On-premise PostgreSQL + pgvector VECTOR(768), tsvector full-text + GIN index, private local disk"),
        ("AI Gateway + LLM", "OmniRoute local gateway (port 20128) → Qwen3:8b; Qwen3-VL:4b vision; CPU fallback via Ollama"),
        ("Retrieval & Reranking", "Hybrid RRF (pgvector cosine + tsvector), FlashRank ms-marco-TinyBERT-L-2-v2 cross-encoder on CPU"),
        ("Ingestion & Parsing", "PyMuPDF + pypdf, python-docx, openpyxl, python-pptx, 13 formats, on-premise local file storage")
    ]
    coords_stk = [
        (0.6, 1.6), (5.1, 1.6),
        (0.6, 2.8), (5.1, 2.8),
        (0.6, 4.0), (5.1, 4.0)
    ]
    for i, (title, desc) in enumerate(stacks):
        x, y = coords_stk[i]
        card = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(1.05))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER

        tb_s = s14.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(4.0), Inches(0.85))
        tf_s = tb_s.text_frame
        tf_s.word_wrap = True

        p1 = tf_s.paragraphs[0]
        r1 = p1.add_run()
        r1.text = title
        r1.font.name = FONT_NAME
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = BRAND_RED

        p2 = tf_s.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = desc
        r2.font.name = FONT_NAME
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = COLOR_BODY

    add_footer(s14, 14)

    # =========================================================================
    # SLIDE 15 — FINAL VALUE PROPOSITION (SOVEREIGNTY STATEMENT)
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    add_header(s15, "CASE STUDY 06.", "CONCLUSION")

    card_fin = s15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.85), Inches(8.8), Inches(4.15))
    card_fin.fill.solid()
    card_fin.fill.fore_color.rgb = COLOR_CARD_BG
    card_fin.line.color.rgb = BRAND_RED

    tb_fin = s15.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(8.0), Inches(3.5))
    tf_fin = tb_fin.text_frame
    tf_fin.word_wrap = True

    p = tf_fin.paragraphs[0]
    r = p.add_run()
    r.text = "Complete Data Sovereignty. Total Privacy. Zero Cloud Dependency."
    r.font.name = FONT_NAME
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE

    p2 = tf_fin.add_paragraph()
    p2.space_before = Pt(14)
    r2 = p2.add_run()
    r2.text = "SipraHub Fully Local RAG empowers enterprises to harness conversational intelligence with complete confidence that proprietary documents never leave private infrastructure."
    r2.font.name = FONT_NAME
    r2.font.size = Pt(13)
    r2.font.color.rgb = COLOR_BODY

    bullets = [
        "100% Private Execution: All LLM inference runs through local OmniRoute gateway — zero external token billing.",
        "Agentic 10-Route Intelligence: Rule-based intent router dispatches DOCUMENT_QA, WEB, CALCULATOR, HYBRID and 6 more routes.",
        "Verifiable Source Citations: Every answer cites chunk_id, document_version_id, similarity score, and section title.",
        "Enterprise-Grade Security: JWT HS256 auth, prompt injection neutralization, bcrypt hashing, zero telemetry."
    ]
    for b in bullets:
        pb = tf_fin.add_paragraph()
        pb.space_before = Pt(8)
        rb = pb.add_run()
        rb.text = f"- {b}"
        rb.font.name = FONT_NAME
        rb.font.size = Pt(11)
        rb.font.color.rgb = COLOR_BODY

    p_close = tf_fin.add_paragraph()
    p_close.space_before = Pt(24)
    rc1 = p_close.add_run()
    rc1.text = "SipraHub  |  AI-Driven Digital Transformation  |  "
    rc1.font.name = FONT_NAME
    rc1.font.size = Pt(11)
    rc1.font.bold = True
    rc1.font.color.rgb = COLOR_TITLE

    rc2 = p_close.add_run()
    rc2.text = "www.siprahub.com"
    rc2.font.name = FONT_NAME
    rc2.font.size = Pt(11)
    rc2.font.bold = True
    rc2.font.color.rgb = BRAND_RED

    add_footer(s15, 15)

    prs.save(OUTPUT_PATH)
    print(f"Successfully generated 100% Fully Local RAG presentation at: {OUTPUT_PATH}")

if __name__ == "__main__":
    create_deck()
